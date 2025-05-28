# presentation_pipeline.py

from langchain_openai import ChatOpenAI
from langchain.prompts import PromptTemplate
import google.generativeai as genai
import os
from typing import Dict, Any, List, Tuple
from pathlib import Path
import base64
from io import BytesIO
import tempfile

# Presentation handling libraries
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import fitz  # PyMuPDF for PDF presentations
from pdf2image import convert_from_path
import json

class PresentationPipeline:
    def __init__(self, api_key: str = None, google_api_key: str = None):
        """Initialize the presentation extraction pipeline with OpenAI and Google Gemini"""
        self.llm = ChatOpenAI(
            model="gpt-4",
            temperature=0.3,
            api_key=api_key or os.getenv("OPENAI_API_KEY")
        )
        
        # Configure Google Gemini for vision capabilities
        genai.configure(api_key=google_api_key or os.getenv("GOOGLE_API_KEY"))
        self.vision_model = genai.GenerativeModel('gemini-1.5-flash')
        
        self.text_extraction_prompt = PromptTemplate(
            input_variables=["slides_content", "presentation_title"],
            template="""
You are an expert at analyzing presentations and extracting key insights for LinkedIn blog posts.

Presentation Title: {presentation_title}

Slides Content:
{slides_content}

Please extract and organize the following information:

1. **Main Topic & Core Message**: What is the presentation about? What's the key takeaway?
2. **Key Points by Section**: Organize the main points logically
3. **Data & Statistics**: Any important numbers, metrics, or data points
4. **Unique Insights**: What makes this presentation valuable or different?
5. **Actionable Takeaways**: What can the audience learn or apply?
6. **Notable Quotes**: Any powerful statements or memorable phrases
7. **Visual Elements Summary**: Description of key charts, diagrams, or visual concepts
8. **Target Audience**: Who would benefit most from this content?
9. **Professional Applications**: How can this be applied in a business context?
10. **Story Arc**: The narrative flow of the presentation

Format the output as a structured analysis suitable for creating an engaging LinkedIn blog post.
Focus on extracting value that resonates with professionals.
"""
        )
        
        self.image_analysis_prompt = """
Analyze this image from a presentation slide and provide:
1. What type of visual is this? (chart, diagram, infographic, photo, etc.)
2. Key information or data shown
3. Main message or insight from this visual
4. How this supports the presentation's narrative
5. Any text or labels visible in the image

Focus on extracting insights that would be valuable for a LinkedIn blog post.
"""
        
        self.supported_formats = ['.pptx', '.ppt', '.pdf', '.odp']
    
    def extract_from_pptx(self, file_path: str) -> Tuple[List[Dict], List[Tuple]]:
        """Extract text and images from PowerPoint presentation"""
        prs = Presentation(file_path)
        slides_data = []
        images_data = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = {
                "slide_number": slide_num,
                "title": "",
                "content": [],
                "speaker_notes": ""
            }
            
            # Extract title
            if slide.shapes.title:
                slide_content["title"] = slide.shapes.title.text
            
            # Extract text and images from shapes
            for shape in slide.shapes:
                # Extract text
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text and text != slide_content["title"]:
                        slide_content["content"].append(text)
                
                # Extract images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        image_format = image.ext
                        
                        # Convert to PIL Image
                        pil_image = Image.open(BytesIO(image_bytes))
                        
                        # Store image data
                        images_data.append((
                            slide_num,
                            pil_image,
                            f"slide_{slide_num}_image_{len(images_data)+1}.{image_format}"
                        ))
                    except Exception as e:
                        print(f"Error extracting image from slide {slide_num}: {e}")
                
                # Extract charts/tables (as text representation)
                if shape.has_chart:
                    slide_content["content"].append("[Chart present - data visualization]")
                elif shape.has_table:
                    table_text = self._extract_table_text(shape.table)
                    slide_content["content"].append(f"[Table]: {table_text}")
            
            # Extract speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text:
                    slide_content["speaker_notes"] = notes_text
            
            slides_data.append(slide_content)
        
        return slides_data, images_data
    
    def extract_from_pdf(self, file_path: str) -> Tuple[List[Dict], List[Tuple]]:
        """Extract content from PDF presentations"""
        slides_data = []
        images_data = []
        
        # Open PDF
        pdf_document = fitz.open(file_path)
        
        for page_num, page in enumerate(pdf_document, 1):
            slide_content = {
                "slide_number": page_num,
                "title": f"Slide {page_num}",
                "content": [],
                "speaker_notes": ""
            }
            
            # Extract text
            text = page.get_text()
            if text:
                lines = text.strip().split('\n')
                # Assume first line might be title
                if lines:
                    slide_content["title"] = lines[0]
                    slide_content["content"] = lines[1:]
            
            # Extract images
            image_list = page.get_images()
            for img_index, img in enumerate(image_list):
                try:
                    # Get image data
                    xref = img[0]
                    pix = fitz.Pixmap(pdf_document, xref)
                    
                    # Convert to PIL Image
                    img_data = pix.tobytes("png")
                    pil_image = Image.open(BytesIO(img_data))
                    
                    images_data.append((
                        page_num,
                        pil_image,
                        f"slide_{page_num}_image_{img_index+1}.png"
                    ))
                except Exception as e:
                    print(f"Error extracting image from slide {page_num}: {e}")
            
            slides_data.append(slide_content)
        
        pdf_document.close()
        return slides_data, images_data
    
    def _extract_table_text(self, table) -> str:
        """Extract text from PowerPoint table"""
        table_text = []
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                row_text.append(cell.text.strip())
            table_text.append(" | ".join(row_text))
        return "\n".join(table_text)
    
    def analyze_image_with_vision(self, image: Image.Image, slide_num: int) -> str:
        """Analyze image using Google Gemini Flash 1.5"""
        try:
            # Prepare the prompt for Gemini
            prompt = f"""
Analyze this image from slide {slide_num} of a presentation and provide:
1. What type of visual is this? (chart, diagram, infographic, photo, etc.)
2. Key information or data shown
3. Main message or insight from this visual
4. How this supports the presentation's narrative
5. Any text or labels visible in the image

Focus on extracting insights that would be valuable for a LinkedIn blog post.
Be specific about data points, trends, or key visual elements that tell a story.
"""
            
            # Generate analysis using Gemini Flash 1.5
            response = self.vision_model.generate_content([prompt, image])
            
            return f"Slide {slide_num} Visual Analysis (Gemini Flash 1.5):\n{response.text}"
            
        except Exception as e:
            print(f"Error analyzing image with Gemini Flash 1.5: {e}")
            return f"Slide {slide_num}: [Visual element present but not analyzed - {str(e)}]"
    
    def extract_from_presentation(self, file_path: str, analyze_images: bool = True) -> Dict[str, Any]:
        """Extract information from any supported presentation format"""
        try:
            file_ext = Path(file_path).suffix.lower()
            
            if file_ext not in self.supported_formats:
                return {
                    "error": f"Unsupported format. Supported: {', '.join(self.supported_formats)}",
                    "status": "error"
                }
            
            # Extract content based on format
            if file_ext in ['.pptx', '.ppt']:
                slides_data, images_data = self.extract_from_pptx(file_path)
            elif file_ext == '.pdf':
                slides_data, images_data = self.extract_from_pdf(file_path)
            else:
                # For other formats, convert to PDF first
                # This would require additional libraries like unoconv
                return {
                    "error": f"Format {file_ext} support coming soon",
                    "status": "error"
                }
            
            # Format slides content
            formatted_slides = self._format_slides_content(slides_data)
            
            # Analyze images if requested
            image_analyses = []
            if analyze_images and images_data:
                print(f"Analyzing {len(images_data)} images...")
                for slide_num, image, filename in images_data[:5]:  # Limit to 5 images
                    analysis = self.analyze_image_with_vision(image, slide_num)
                    image_analyses.append(analysis)
            
            # Combine all content
            presentation_title = Path(file_path).stem.replace('_', ' ').title()
            
            # Get main analysis using LLM
            full_content = formatted_slides
            if image_analyses:
                full_content += "\n\n=== Visual Elements Analysis ===\n"
                full_content += "\n\n".join(image_analyses)
            
            response = self.llm.invoke(
                self.text_extraction_prompt.format(
                    slides_content=full_content,
                    presentation_title=presentation_title
                )
            )
            
            return {
                "source_type": "presentation",
                "file_path": file_path,
                "title": presentation_title,
                "total_slides": len(slides_data),
                "images_found": len(images_data),
                "images_analyzed": len(image_analyses),
                "extracted_info": response.content,
                "status": "success"
            }
            
        except Exception as e:
            return {
                "source_type": "presentation",
                "file_path": file_path,
                "error": str(e),
                "status": "error"
            }
    
    def _format_slides_content(self, slides_data: List[Dict]) -> str:
        """Format slides data into readable text"""
        formatted_content = []
        
        for slide in slides_data:
            slide_text = f"\n=== Slide {slide['slide_number']}: {slide['title']} ===\n"
            
            if slide['content']:
                slide_text += "Content:\n"
                for content in slide['content']:
                    slide_text += f"- {content}\n"
            
            if slide['speaker_notes']:
                slide_text += f"\nSpeaker Notes:\n{slide['speaker_notes']}\n"
            
            formatted_content.append(slide_text)
        
        return "\n".join(formatted_content)
    
    def extract_key_slides(self, file_path: str, key_slide_numbers: List[int]) -> Dict[str, Any]:
        """Extract only specific slides from presentation"""
        try:
            # Extract all slides first
            if file_path.lower().endswith(('.pptx', '.ppt')):
                slides_data, images_data = self.extract_from_pptx(file_path)
            else:
                slides_data, images_data = self.extract_from_pdf(file_path)
            
            # Filter to only key slides
            key_slides = [s for s in slides_data if s['slide_number'] in key_slide_numbers]
            key_images = [(n, img, f) for n, img, f in images_data if n in key_slide_numbers]
            
            # Process only key slides
            formatted_slides = self._format_slides_content(key_slides)
            
            return {
                "source_type": "presentation_excerpt",
                "file_path": file_path,
                "key_slides": key_slide_numbers,
                "extracted_info": formatted_slides,
                "status": "success"
            }
            
        except Exception as e:
            return {
                "source_type": "presentation_excerpt",
                "file_path": file_path,
                "error": str(e),
                "status": "error"
            }