import time
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import Dict, Any, List
from ingestion.config import ExtractedContent, ContentType, ProcessingModel, DocumentMetadata

class PPTProcessor:
    """Process PowerPoint files and extract content"""
    
    @staticmethod
    def extract_content(file_path: str) -> ExtractedContent:
        """Extract text and slide content from PowerPoint"""
        start_time = time.time()
        
        # Load PowerPoint presentation
        presentation = Presentation(file_path)
        
        raw_text = ""
        slides_data = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_content = {
                "slide_number": slide_num,
                "title": "",
                "content": [],
                "notes": "",
                "shapes_count": len(slide.shapes),
                "has_images": False,
                "images": []
            }
            
            # Extract text from shapes
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    slide_text += text + "\n"
                    
                    # Check if it's likely a title (first text shape or larger font)
                    if not slide_content["title"] and len(text) < 100:
                        slide_content["title"] = text
                    else:
                        slide_content["content"].append(text)
                
                # Check for images and extract bytes
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and hasattr(shape, "image"):
                        image = shape.image
                        image_bytes = image.blob
                        image_ext = image.ext
                        mime_map = {
                            "png": "image/png",
                            "jpg": "image/jpeg",
                            "jpeg": "image/jpeg",
                            "gif": "image/gif",
                            "bmp": "image/bmp",
                            "tiff": "image/tiff",
                            "wmf": "image/wmf",
                        }
                        mime_type = mime_map.get(image_ext.lower(), "image/png")
                        slide_content["images"].append({
                            "mime_type": mime_type,
                            "bytes_len": len(image_bytes),
                            "image_bytes": image_bytes
                        })
                        slide_content["has_images"] = True
                except Exception:
                    pass
            
            # Extract notes
            if slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_content["notes"] = notes_text
                    slide_text += f"\nNotes: {notes_text}\n"
            
            slides_data.append(slide_content)
            raw_text += f"Slide {slide_num}:\n{slide_text}\n{'='*50}\n"
        
        # Structure the extracted data
        structured_data = {
            "total_slides": len(presentation.slides),
            "slides": slides_data,
            "presentation_metadata": {
                "has_master_slides": len(presentation.slide_masters) > 0,
                "total_shapes": sum(len(slide.shapes) for slide in presentation.slides),
                "slides_with_images": sum(1 for slide_data in slides_data if slide_data["has_images"])
            }
        }
        
        # Additional metadata
        metadata = {
            "total_slides": len(presentation.slides),
            "total_characters": len(raw_text),
            "word_count": len(raw_text.split()),
            "slides_with_content": len([s for s in slides_data if s["content"] or s["title"]]),
            "slides_with_notes": len([s for s in slides_data if s["notes"]])
        }
        
        processing_time = time.time() - start_time
        
        return ExtractedContent(
            content_type=ContentType.POWERPOINT,
            file_path=file_path,
            raw_text=raw_text.strip(),
            structured_data=structured_data,
            metadata=metadata,
            processing_model=ProcessingModel.GROQ_LLAMA_8B,
            processing_time=processing_time
        )
    
    @staticmethod
    def update_document_metadata(metadata: DocumentMetadata, extracted: ExtractedContent) -> DocumentMetadata:
        """Update document metadata with PowerPoint-specific information"""
        metadata.page_count = extracted.metadata.get("total_slides", 0)
        metadata.word_count = extracted.metadata.get("word_count", 0)
        
        return metadata