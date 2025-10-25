"""All format-specific content extraction handlers"""

import os
import sys
from pathlib import Path
from typing import Dict, Any, Optional
import base64

# Add parent directory to path for langsmith_config import
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from langsmith_config import trace_step

from .config import ContentType, ExtractedContent


class FormatHandler:
    """Unified handler for all file formats using appropriate loaders"""
    
    @staticmethod
    @trace_step("extract_pdf", "tool")
    def extract_pdf(file_path: str) -> ExtractedContent:
        """Extract content from PDF using document loader"""
        try:
            from langchain_community.document_loaders import PyPDFLoader
            
            loader = PyPDFLoader(file_path)
            pages = loader.load()
            
            # Combine all pages
            full_content = "\n\n".join([page.page_content for page in pages])
            
            # Extract metadata from first page
            metadata = {
                "page_count": len(pages),
                "source": file_path,
                **(pages[0].metadata if pages else {})
            }
            
            # Extract structured data (page-by-page breakdown)
            structured_data = {
                "pages": [
                    {
                        "page_num": i + 1,
                        "content": page.page_content,
                        "metadata": page.metadata
                    }
                    for i, page in enumerate(pages)
                ]
            }
            
            return ExtractedContent(
                content=full_content,
                metadata=metadata,
                structured_data=structured_data
            )
            
        except Exception as e:
            raise RuntimeError(f"PDF extraction failed: {str(e)}")
    
    @staticmethod
    @trace_step("extract_word", "tool")
    def extract_word(file_path: str) -> ExtractedContent:
        """Extract content from Word documents"""
        try:
            from langchain_community.document_loaders import Docx2txtLoader
            
            loader = Docx2txtLoader(file_path)
            documents = loader.load()
            
            # Combine all document parts
            content = "\n\n".join([doc.page_content for doc in documents])
            
            metadata = {
                "source": file_path,
                "document_count": len(documents),
                **(documents[0].metadata if documents else {})
            }
            
            return ExtractedContent(
                content=content,
                metadata=metadata
            )
            
        except Exception as e:
            raise RuntimeError(f"Word extraction failed: {str(e)}")
    
    @staticmethod
    @trace_step("extract_powerpoint", "tool")
    def extract_powerpoint(file_path: str) -> ExtractedContent:
        """Extract content from PowerPoint presentations"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            
            slides_content = []
            slides_data = []
            
            for idx, slide in enumerate(prs.slides, 1):
                # Extract text from all shapes
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                combined_text = "\n".join(slide_text)
                slides_content.append(f"Slide {idx}:\n{combined_text}")
                
                slides_data.append({
                    "slide_number": idx,
                    "content": combined_text,
                    "shape_count": len(slide.shapes)
                })
            
            full_content = "\n\n".join(slides_content)
            
            metadata = {
                "source": file_path,
                "slide_count": len(prs.slides),
                "has_images": any(
                    hasattr(shape, "image") 
                    for slide in prs.slides 
                    for shape in slide.shapes
                )
            }
            
            structured_data = {
                "slides": slides_data
            }
            
            return ExtractedContent(
                content=full_content,
                metadata=metadata,
                structured_data=structured_data
            )
            
        except Exception as e:
            raise RuntimeError(f"PowerPoint extraction failed: {str(e)}")
    
    @staticmethod
    @trace_step("extract_code", "tool")
    def extract_code(file_path: str) -> ExtractedContent:
        """Extract content from code files"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            # Get file extension for language detection
            ext = Path(file_path).suffix.lstrip('.')
            
            # Basic code analysis
            lines = content.split('\n')
            non_empty_lines = [line for line in lines if line.strip()]
            
            # Simple comment detection (works for most languages)
            comment_chars = ['#', '//', '/*', '*', '--']
            comment_lines = sum(
                1 for line in non_empty_lines 
                if any(line.strip().startswith(char) for char in comment_chars)
            )
            
            metadata = {
                "source": file_path,
                "language": ext,
                "total_lines": len(lines),
                "code_lines": len(non_empty_lines),
                "comment_lines": comment_lines,
                "file_size_bytes": len(content)
            }
            
            structured_data = {
                "language": ext,
                "has_imports": any(
                    line.strip().startswith(('import', 'from', 'require', 'using', 'include'))
                    for line in non_empty_lines
                ),
                "has_functions": any(
                    keyword in content.lower() 
                    for keyword in ['def ', 'function ', 'func ', 'fn ', 'sub ']
                ),
                "has_classes": any(
                    keyword in content.lower()
                    for keyword in ['class ', 'struct ', 'interface ']
                )
            }
            
            return ExtractedContent(
                content=content,
                metadata=metadata,
                structured_data=structured_data
            )
            
        except Exception as e:
            raise RuntimeError(f"Code extraction failed: {str(e)}")
    
    @staticmethod
    @trace_step("extract_text", "tool")
    def extract_text(file_path: str) -> ExtractedContent:
        """Extract content from plain text files"""
        try:
            from langchain_community.document_loaders import TextLoader
            
            loader = TextLoader(file_path, encoding='utf-8')
            documents = loader.load()
            
            content = "\n\n".join([doc.page_content for doc in documents])
            
            # Basic text statistics
            words = content.split()
            lines = content.split('\n')
            
            metadata = {
                "source": file_path,
                "word_count": len(words),
                "line_count": len(lines),
                "character_count": len(content),
                **(documents[0].metadata if documents else {})
            }
            
            return ExtractedContent(
                content=content,
                metadata=metadata
            )
            
        except Exception as e:
            raise RuntimeError(f"Text extraction failed: {str(e)}")
    
    @staticmethod
    @trace_step("extract_markdown", "tool")
    def extract_markdown(file_path: str) -> ExtractedContent:
        """Extract content from Markdown files"""
        try:
            from langchain_community.document_loaders import UnstructuredMarkdownLoader
            
            loader = UnstructuredMarkdownLoader(file_path)
            documents = loader.load()
            
            content = "\n\n".join([doc.page_content for doc in documents])
            
            # Count headers and links
            lines = content.split('\n')
            header_count = sum(1 for line in lines if line.strip().startswith('#'))
            link_count = content.count('](')
            
            metadata = {
                "source": file_path,
                "header_count": header_count,
                "link_count": link_count,
                "character_count": len(content),
                **(documents[0].metadata if documents else {})
            }
            
            return ExtractedContent(
                content=content,
                metadata=metadata
            )
            
        except Exception as e:
            # Fallback to plain text if markdown loader fails
            return FormatHandler.extract_text(file_path)
    
    @staticmethod
    @trace_step("extract_image", "tool")
    def extract_image(file_path: str) -> ExtractedContent:
        """Extract content from image files"""
        try:
            with open(file_path, 'rb') as f:
                image_bytes = f.read()
            
            # Get image format
            ext = Path(file_path).suffix.lstrip('.').lower()
            mime_type = f"image/{ext if ext != 'jpg' else 'jpeg'}"
            
            # Encode for AI processing
            image_b64 = base64.b64encode(image_bytes).decode('utf-8')
            
            metadata = {
                "source": file_path,
                "mime_type": mime_type,
                "file_extension": ext,
                "file_size_bytes": len(image_bytes),
                "has_image_data": True
            }
            
            structured_data = {
                "image_base64": image_b64,
                "mime_type": mime_type
            }
            
            return ExtractedContent(
                content="[Image content - requires vision AI analysis]",
                metadata=metadata,
                image_data=image_bytes,
                structured_data=structured_data
            )
            
        except Exception as e:
            raise RuntimeError(f"Image extraction failed: {str(e)}")
    
    # Main dispatcher
    @classmethod
    @trace_step("extract_content", "tool")
    def extract(cls, file_path: str, content_type: ContentType) -> ExtractedContent:
        """Route to appropriate extraction method"""
        
        handlers = {
            ContentType.PDF: cls.extract_pdf,
            ContentType.WORD: cls.extract_word,
            ContentType.POWERPOINT: cls.extract_powerpoint,
            ContentType.CODE: cls.extract_code,
            ContentType.TEXT: cls.extract_text,
            ContentType.MARKDOWN: cls.extract_markdown,
            ContentType.IMAGE: cls.extract_image,
        }
        
        handler = handlers.get(content_type)
        if not handler:
            raise ValueError(f"No handler for content type: {content_type}")
        
        return handler(file_path)


# Export
__all__ = ['FormatHandler']

